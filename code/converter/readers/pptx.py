"""PowerPoint (.pptx) reader using python-pptx. One block per slide."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu  # noqa: F401 (kept for future layout-aware features)

from .common import Block, ImageAsset


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        indent = max(0, para.level or 0)
        text = "".join(r.text or "" for r in para.runs).strip()
        if text:
            lines.append("  " * indent + "- " + text)
    return "\n".join(lines).strip()


def _is_title(shape) -> bool:
    try:
        ph = shape.placeholder_format
    except (AttributeError, ValueError):
        return False
    if ph is None:
        return False
    return ph.idx == 0 or (ph.type is not None and ph.type in (13, 14, 15))  # TITLE / CENTER_TITLE / SUBTITLE


def _slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes = slide.notes_slide.notes_text_frame.text or ""
    return notes.strip()


def _iter_images(slide, slide_num: int, image_mode: str, images: list[ImageAsset]) -> list[str]:
    markers: list[str] = []
    idx = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            idx += 1
            if image_mode == "omit":
                continue
            if image_mode == "placeholder":
                markers.append(f"_[이미지: 슬라이드 {slide_num} #{idx}]_")
                continue
            try:
                image = shape.image
            except (AttributeError, ValueError):
                continue
            ext = (image.ext or "bin").lower()
            name = f"slide{slide_num}_img{idx}.{ext}"
            images.append(ImageAsset(path=f"images/{name}", data=image.blob))
            markers.append(f"![슬라이드 {slide_num} 이미지 {idx}](images/{name})")
    return markers


def read(path: str | Path, image_mode: str = "extract") -> tuple[list[Block], list[ImageAsset]]:
    path = Path(path)
    prs = Presentation(str(path))
    blocks: list[Block] = [
        Block(
            kind="heading",
            md=f"# {path.name}\n\n> PowerPoint 문서 · 슬라이드 {len(prs.slides)}개\n",
            source=path.name,
        )
    ]
    images: list[ImageAsset] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        bodies: list[str] = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            if not title and _is_title(shape):
                title = text.replace("- ", "", 1).splitlines()[0].strip()
            else:
                bodies.append(text)
        image_markers = _iter_images(slide, i, image_mode, images)
        notes = _slide_notes(slide)
        header = f"---\n\n### 슬라이드 {i}" + (f": {title}" if title else "")
        body_md = "\n\n".join(bodies)
        img_md = "\n\n".join(image_markers)
        notes_md = f"> 노트: {notes}" if notes else ""
        md = "\n\n".join(x for x in [header, body_md, img_md, notes_md] if x) + "\n"
        source = f"슬라이드 {i}" + (f": {title}" if title else "")
        blocks.append(Block(kind="slide", md=md, source=source))
    return blocks, images
