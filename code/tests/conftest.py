"""Generate minimal Office/PDF fixtures at test time — no external sample files needed."""

from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture(scope="session")
def fixtures_dir(tmp_path_factory) -> Path:
    return tmp_path_factory.mktemp("fixtures")


@pytest.fixture(scope="session")
def sample_xlsx(fixtures_dir: Path) -> Path:
    from openpyxl import Workbook

    path = fixtures_dir / "sample.xlsx"
    wb = Workbook()
    sheet1 = wb.active
    sheet1.title = "인사"
    sheet1.append(["이름", "부서", "입사년도"])
    sheet1.append(["김철수", "영업", 2020])
    sheet1.append(["이영희", "개발", 2021])
    sheet1.append(["박민수", "기획", 2019])

    sheet2 = wb.create_sheet("매출")
    sheet2.append(["월", "매출"])
    for i in range(1, 13):
        sheet2.append([f"{i}월", i * 1000])

    wb.save(path)
    return path


@pytest.fixture(scope="session")
def sample_docx(fixtures_dir: Path) -> Path:
    from docx import Document

    path = fixtures_dir / "sample.docx"
    doc = Document()
    doc.add_heading("테스트 문서", level=1)
    doc.add_paragraph("첫 번째 단락입니다.")
    doc.add_heading("2장. 소개", level=2)
    doc.add_paragraph("두 번째 단락은 일반 본문입니다.")
    doc.add_paragraph("불릿 항목 1", style="List Bullet")
    doc.add_paragraph("불릿 항목 2", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = "2"
    doc.save(path)
    return path


@pytest.fixture(scope="session")
def sample_pptx(fixtures_dir: Path) -> Path:
    from pptx import Presentation

    path = fixtures_dir / "sample.pptx"
    prs = Presentation()
    for i in range(1, 4):
        layout = prs.slide_layouts[1] if i != 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"슬라이드 {i} 제목"
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            body.text = f"슬라이드 {i} 본문 첫 줄"
            body.text_frame.add_paragraph().text = f"슬라이드 {i} 본문 둘째 줄"
        slide.notes_slide.notes_text_frame.text = f"발표자 노트 {i}"
    prs.save(path)
    return path


@pytest.fixture(scope="session")
def sample_pdf(fixtures_dir: Path) -> Path:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4

    path = fixtures_dir / "sample.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setFont("Helvetica", 12)
    c.drawString(72, 800, "Page 1 line 1")
    c.drawString(72, 780, "Page 1 line 2")
    c.showPage()
    c.drawString(72, 800, "Page 2 line 1")
    c.showPage()
    c.save()
    return path
